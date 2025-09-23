def ler_arquivo(caminho):
    if caminho.endswith(".txt"):
        with open(caminho, "r", encoding="utf-8") as f:
            return f.read()

    elif caminho.endswith(".pdf"):
        import PyPDF2
        texto = ""
        with open(caminho, "rb") as f:
            pdf = PyPDF2.PdfReader(f)
            for pagina in pdf.pages:
                texto += pagina.extract_text() + "\n"
        return texto

    elif caminho.endswith(".docx"):
        import docx
        doc = docx.Document(caminho)
        return "\n".join([p.text for p in doc.paragraphs])

    elif caminho.endswith(".xlsx"):
        from openpyxl import load_workbook
        wb = load_workbook(caminho)
        texto = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                linha = [str(cell) if cell is not None else "" for cell in row]
                texto += " | ".join(linha) + "\n"
        return texto

    elif caminho.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(caminho)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto

    elif caminho.endswith(".ppt"):  # suporte básico para PPT antigo
        return "Arquivos PPT antigos não são totalmente suportados. Converta para PPTX."

    else:
        return "Formato de arquivo não suportado"
